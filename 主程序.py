#!/usr/bin/env python3
"""
单词PPT生成器（轻量级版本）

该脚本从Excel文件读取单词数据，生成包含单词、音标、词根词缀、例句和释义的PowerPoint演示文稿。
使用openpyxl直接读取Excel文件，不依赖pandas，减小打包体积。
添加了Tkinter图形界面，支持文件选择、表格选择、预览和生成功能。
"""

import sys
import argparse
import os
from typing import Dict, List, Optional
from openpyxl import load_workbook

# 尝试导入Tkinter，支持图形界面和命令行两种模式
try:
    import tkinter as tk
    from tkinter import filedialog, ttk, messagebox
    from tkinter.scrolledtext import ScrolledText
    HAS_GUI = True
except ImportError:
    HAS_GUI = False

# 依赖检查已移除，因为在打包时所有依赖都已包含在可执行文件中

from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from pptx import Presentation

# 常量定义
SLIDE_WIDTH = Inches(16)
SLIDE_HEIGHT = Inches(9)
WORD_FONT_SIZE = Pt(72)
PHONETIC_FONT_SIZE = Pt(32)
CONTENT_FONT_SIZE = Pt(32)
TEXT_WRAP_THRESHOLD = 40

# 列名映射
REQUIRED_COLUMNS = {
    '英文单词': 'word',
    '英文音标': 'phonetic',
    '词根词缀': 'morphology',
    '例句': 'example',
    '例句释义': 'example_def',
    '单词释义': 'word_def'
}

def process_text(text: str, threshold: int = TEXT_WRAP_THRESHOLD) -> str:
    """处理文本，超过阈值时自动换行"""
    if len(text) <= threshold:
        return text
    
    lines = [text[i:i+threshold] for i in range(0, len(text), threshold)]
    return '\n                  '.join(lines)

def load_data(input_file: str) -> List[Dict[str, str]]:
    """从Excel文件加载数据并进行验证"""
    try:
        wb = load_workbook(input_file)
        ws = wb.active
    except Exception as e:
        print(f"读取Excel文件时发生错误：{e}")
        sys.exit(1)
    
    # 获取表头
    headers = [cell.value for cell in ws[1]]
    
    # 验证必要的列是否存在
    missing_columns = [col for col in REQUIRED_COLUMNS if col not in headers]
    if missing_columns:
        print(f"Excel文件缺少必要的列：{', '.join(missing_columns)}")
        sys.exit(1)
    
    # 加载数据
    data = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        row_dict = {}
        for i, value in enumerate(row):
            if i < len(headers):
                row_dict[headers[i]] = str(value) if value is not None else ''
        data.append(row_dict)
    
    return data

def create_presentation() -> Presentation:
    """创建新的演示文稿并设置页面大小"""
    ppt = Presentation()
    ppt.slide_width = SLIDE_WIDTH
    ppt.slide_height = SLIDE_HEIGHT
    return ppt

def add_word_textbox(slide, left: Inches, top: Inches, width: Inches, height: Inches, word: str) -> None:
    """添加单词文本框"""
    word_textbox = slide.shapes.add_textbox(left, top, width, height).text_frame
    word_paragraph = word_textbox.add_paragraph()
    word_paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    word_paragraph.font.size = WORD_FONT_SIZE
    word_paragraph.text = word
    word_textbox.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

def add_phonetic_textbox(slide, left: Inches, top: Inches, width: Inches, height: Inches, phonetic: str) -> None:
    """添加音标文本框"""
    phonetic_textbox = slide.shapes.add_textbox(left, top, width, height).text_frame
    phonetic_paragraph = phonetic_textbox.add_paragraph()
    phonetic_paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    phonetic_paragraph.font.size = PHONETIC_FONT_SIZE
    phonetic_paragraph.text = str(phonetic)
    phonetic_textbox.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

def add_morphology_textbox(slide, left: Inches, top: Inches, width: Inches, height: Inches, morphology: str) -> None:
    """添加词根词缀文本框"""
    morphology_textbox = slide.shapes.add_textbox(left, top, width, height).text_frame
    morphology_textbox.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    morphology_textbox.word_wrap = True
    
    morphology_paragraph = morphology_textbox.add_paragraph()
    morphology_paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    morphology_paragraph.font.size = CONTENT_FONT_SIZE
    
    processed_morphology = process_text(str(morphology))
    morphology_paragraph.text = "词根词缀：" + processed_morphology

def add_example_textbox(slide, left: Inches, top: Inches, width: Inches, height: Inches, 
                       example: str, example_def: str) -> None:
    """添加例句和例句释义文本框"""
    example_textbox = slide.shapes.add_textbox(left, top, width, height).text_frame
    example_textbox.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    example_textbox.word_wrap = True
    
    example_paragraph = example_textbox.add_paragraph()
    example_paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    example_paragraph.font.size = CONTENT_FONT_SIZE
    
    processed_example = process_text(str(example))
    processed_example_def = process_text(str(example_def))
    example_paragraph.text = f"例句：{processed_example}\n例句释义：{processed_example_def}"

def add_word_def_textbox(slide, left: Inches, top: Inches, width: Inches, height: Inches, word_def: str) -> None:
    """添加单词释义文本框"""
    word_def_textbox = slide.shapes.add_textbox(left, top, width, height).text_frame
    word_def_textbox.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    
    word_def_paragraph = word_def_textbox.add_paragraph()
    word_def_paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    word_def_paragraph.font.size = CONTENT_FONT_SIZE
    
    processed_word_def = process_text(str(word_def))
    word_def_paragraph.text = "单词释义：" + processed_word_def

def generate_slide(ppt: Presentation, row: Dict[str, str]) -> None:
    """为单个单词生成幻灯片"""
    # 添加幻灯片布局（空白布局）
    slide_layout = ppt.slide_layouts[6]
    slide = ppt.slides.add_slide(slide_layout)
    
    # 单词文本框位置和大小
    word_left, word_top = Inches(4), Inches(1)
    word_width, word_height = Inches(8), Inches(2)
    
    # 添加单词文本框
    add_word_textbox(slide, word_left, word_top, word_width, word_height, row['英文单词'])
    
    # 添加音标文本框
    phonetic_left = word_left + Inches(3)
    phonetic_top = word_top + Inches(1.4)
    phonetic_width, phonetic_height = Inches(2), Inches(1)
    add_phonetic_textbox(slide, phonetic_left, phonetic_top, phonetic_width, phonetic_height, row['英文音标'])
    
    # 添加词根词缀文本框
    morphology_left, morphology_top = Inches(2), word_top + Inches(3)
    morphology_width, morphology_height = Inches(12), Inches(2)
    add_morphology_textbox(slide, morphology_left, morphology_top, morphology_width, morphology_height, row['词根词缀'])
    
    # 添加例句和释义文本框
    example_left, example_top = Inches(2), word_top + Inches(4.4)
    example_width, example_height = Inches(12), Inches(3)
    add_example_textbox(slide, example_left, example_top, example_width, example_height, 
                       row['例句'], row['例句释义'])
    
    # 添加单词释义文本框（如果存在）
    if row['单词释义']:
        word_def_left, word_def_top = Inches(2), word_top + Inches(2)
        word_def_width, word_def_height = Inches(8), Inches(1.5)
        add_word_def_textbox(slide, word_def_left, word_def_top, word_def_width, word_def_height, row['单词释义'])

def generate_ppt(input_file: str, output_file: str) -> int:
    """生成PPT文件并返回处理的单词数量"""
    # 使用默认值
    if input_file is None:
        input_file = 'words.xlsx'
    if output_file is None:
        output_file = 'words.pptx'
    
    # 加载数据
    data = load_data(input_file)
    
    # 创建演示文稿
    ppt = create_presentation()
    
    # 生成幻灯片
    word_count = len(data)
    print(f"开始生成PPT，共 {word_count} 个单词...")
    
    for index, row in enumerate(data):
        if (index + 1) % 10 == 0 or index + 1 == word_count:
            print(f"处理进度：{index + 1}/{word_count}")
        generate_slide(ppt, row)
    
    # 保存演示文稿
    try:
        ppt.save(output_file)
        print(f"\n成功生成PPT文件：{output_file}")
        print(f"共处理 {word_count} 个单词")
        return word_count
    except Exception as e:
        print(f'保存演示文稿时发生错误：{e}')
        return 0

def parse_args() -> argparse.Namespace:
    """解析命令行参数"""
    parser = argparse.ArgumentParser(description='单词PPT生成器（轻量级版本）')
    parser.add_argument('-i', '--input', default=None, help='输入Excel文件路径')
    parser.add_argument('-o', '--output', default=None, help='输出PPT文件路径')
    parser.add_argument('--version', action='version', version='%(prog)s 1.0.0')
    return parser.parse_args()

class PPTGeneratorGUI:
    """PPT生成器图形界面"""
    
    def __init__(self, master):
        """初始化图形界面"""
        self.master = master
        master.title("单词PPT生成器")
        master.geometry("800x550")  # 增加默认窗口高度，确保所有按钮可见
        master.resizable(True, True)
        master.minsize(800, 550)  # 设置最小窗口大小，确保所有控件都能显示
        
        # 设置主题
        self.style = ttk.Style()
        try:
            self.style.theme_use('clam')
        except:
            pass
        
        # 创建主框架
        self.main_frame = ttk.Frame(master, padding="10")
        self.main_frame.pack(fill=tk.BOTH, expand=True)
        
        # 创建输入文件选择区域
        self.input_frame = ttk.LabelFrame(self.main_frame, text="输入文件", padding="10")
        self.input_frame.pack(fill=tk.X, pady=5)
        
        self.input_path_var = tk.StringVar()
        self.input_entry = ttk.Entry(self.input_frame, textvariable=self.input_path_var, width=60)
        self.input_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.browse_button = ttk.Button(self.input_frame, text="浏览", command=self.browse_input_file)
        self.browse_button.pack(side=tk.RIGHT, padx=5)
        
        # 创建表格选择区域
        self.sheet_frame = ttk.LabelFrame(self.main_frame, text="表格选择", padding="10")
        self.sheet_frame.pack(fill=tk.X, pady=5)
        
        self.sheet_var = tk.StringVar()
        self.sheet_combo = ttk.Combobox(self.sheet_frame, textvariable=self.sheet_var, width=50)
        self.sheet_combo.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        self.sheet_combo['state'] = 'disabled'
        
        self.refresh_sheets_button = ttk.Button(self.sheet_frame, text="刷新表格", command=self.refresh_sheets)
        self.refresh_sheets_button.pack(side=tk.RIGHT, padx=5)
        self.refresh_sheets_button['state'] = 'disabled'
        
        # 创建输出文件选择区域
        self.output_frame = ttk.LabelFrame(self.main_frame, text="输出文件", padding="10")
        self.output_frame.pack(fill=tk.X, pady=5)
        
        self.output_path_var = tk.StringVar()
        self.output_path_var.set("words.pptx")
        self.output_entry = ttk.Entry(self.output_frame, textvariable=self.output_path_var, width=60)
        self.output_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.output_browse_button = ttk.Button(self.output_frame, text="浏览", command=self.browse_output_file)
        self.output_browse_button.pack(side=tk.RIGHT, padx=5)
        
        # 创建模板区域
        self.template_frame = ttk.LabelFrame(self.main_frame, text="表格模板", padding="10")
        self.template_frame.pack(fill=tk.X, pady=5)
        
        self.template_info = ttk.Label(self.template_frame, text="点击下方按钮打开示例表格模板，了解正确的表格格式")
        self.template_info.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        self.template_button = ttk.Button(self.template_frame, text="打开模板", command=self.open_template)
        self.template_button.pack(side=tk.RIGHT, padx=5)
        
        # 创建按钮区域
        self.button_frame = ttk.Frame(self.main_frame, padding="10")
        self.button_frame.pack(fill=tk.X, pady=10)  # 只横向填充，确保按钮区域始终可见
        
        self.generate_button = ttk.Button(self.button_frame, text="生成PPT", command=self.generate)
        self.generate_button.pack(fill=tk.X, padx=5, pady=5)  # 横向填充，确保按钮完整显示
        
        # 创建状态栏
        self.status_var = tk.StringVar()
        self.status_var.set("就绪")
        self.status_bar = ttk.Label(master, textvariable=self.status_var, relief=tk.SUNKEN, anchor=tk.W)
        self.status_bar.pack(side=tk.BOTTOM, fill=tk.X)
        
        # 存储当前选择的Excel文件和表格
        self.current_workbook = None
        self.current_sheet = None
    
    def browse_input_file(self):
        """浏览选择输入Excel文件"""
        file_path = filedialog.askopenfilename(
            title="选择Excel文件",
            filetypes=[("Excel文件", "*.xlsx *.xls"), ("所有文件", "*.*")]
        )
        if file_path:
            self.input_path_var.set(file_path)
            self.status_var.set(f"已选择文件: {os.path.basename(file_path)}")
            
            # 尝试加载文件并获取表格列表
            try:
                self.current_workbook = load_workbook(file_path)
                sheet_names = self.current_workbook.sheetnames
                self.sheet_combo['values'] = sheet_names
                if sheet_names:
                    self.sheet_var.set(sheet_names[0])
                    self.current_sheet = sheet_names[0]
                self.sheet_combo['state'] = 'normal'
                self.refresh_sheets_button['state'] = 'normal'
                self.status_var.set(f"已加载文件: {os.path.basename(file_path)}，包含 {len(sheet_names)} 个表格")
            except Exception as e:
                messagebox.showerror("错误", f"加载文件失败: {str(e)}")
                self.status_var.set("加载文件失败")
                self.sheet_combo['state'] = 'disabled'
                self.refresh_sheets_button['state'] = 'disabled'
    
    def browse_output_file(self):
        """浏览选择输出PPT文件"""
        file_path = filedialog.asksaveasfilename(
            title="选择输出PPT文件",
            defaultextension=".pptx",
            filetypes=[("PowerPoint文件", "*.pptx"), ("所有文件", "*.*")]
        )
        if file_path:
            self.output_path_var.set(file_path)
            self.status_var.set(f"已设置输出文件: {os.path.basename(file_path)}")
    
    def refresh_sheets(self):
        """刷新表格列表"""
        input_path = self.input_path_var.get()
        if not input_path:
            messagebox.showwarning("警告", "请先选择输入文件")
            return
        
        try:
            self.current_workbook = load_workbook(input_path)
            sheet_names = self.current_workbook.sheetnames
            self.sheet_combo['values'] = sheet_names
            if sheet_names:
                self.sheet_var.set(sheet_names[0])
                self.current_sheet = sheet_names[0]
            self.status_var.set(f"已刷新表格列表，包含 {len(sheet_names)} 个表格")
        except Exception as e:
            messagebox.showerror("错误", f"刷新表格失败: {str(e)}")
            self.status_var.set("刷新表格失败")
    

    
    def generate(self):
        """生成PPT文件"""
        input_path = self.input_path_var.get()
        if not input_path:
            messagebox.showwarning("警告", "请先选择输入文件")
            return
        
        sheet_name = self.sheet_var.get()
        if not sheet_name:
            messagebox.showwarning("警告", "请选择表格")
            return
        
        output_path = self.output_path_var.get()
        if not output_path:
            messagebox.showwarning("警告", "请设置输出文件路径")
            return
        
        try:
            self.status_var.set("正在生成PPT...")
            self.master.update()
            
            # 生成PPT
            word_count = generate_ppt_from_sheet(input_path, output_path, sheet_name)
            
            if word_count > 0:
                messagebox.showinfo("成功", f"PPT生成成功！\n共处理 {word_count} 个单词\n输出文件: {output_path}")
                self.status_var.set(f"PPT生成成功，共处理 {word_count} 个单词")
            else:
                messagebox.showwarning("警告", "PPT生成失败，可能是因为表格中没有数据")
                self.status_var.set("PPT生成失败")
        except Exception as e:
            messagebox.showerror("错误", f"生成PPT失败: {str(e)}")
            self.status_var.set(f"生成PPT失败: {str(e)}")
    
    def open_template(self):
        """打开表格模板"""
        import os
        import subprocess
        import sys
        import tempfile
        import time
        
        # 使用中文名作为模板文件名
        template_filename = "单词表模板.xlsx"
        
        # 确定模板文件路径
        if getattr(sys, 'frozen', False):
            # 打包后的环境
            # 对于打包环境，我们优先在可执行文件所在目录创建模板文件
            template_path = os.path.join(os.path.dirname(sys.executable), template_filename)
        else:
            # 开发环境
            template_path = os.path.join(os.path.dirname(__file__), template_filename)
        
        # 检查是否存在原始模板文件
        if os.path.exists(template_path):
            # 使用临时文件来预览，避免修改原始文件
            try:
                # 创建临时目录
                temp_dir = tempfile.gettempdir()
                # 生成临时文件名，避免命名冲突
                temp_template_name = f"单词表模板_{int(time.time())}.xlsx"
                temp_template_path = os.path.join(temp_dir, temp_template_name)
                
                # 复制原始模板文件到临时文件
                import shutil
                shutil.copy2(template_path, temp_template_path)
                
                # 使用默认程序打开临时模板文件，并等待程序关闭
                process = subprocess.Popen(['start', '/wait', temp_template_path], shell=True)
                # 等待进程结束（即用户关闭表格文件）
                process.wait()
                
                # 关闭后立即删除临时文件
                if os.path.exists(temp_template_path):
                    try:
                        os.remove(temp_template_path)
                    except:
                        pass
                
                self.status_var.set("已打开表格模板")
                
            except Exception as e:
                messagebox.showerror("错误", f"打开模板失败: {str(e)}")
                self.status_var.set("打开模板失败")
        else:
            # 如果模板文件不存在，尝试从当前工作目录查找
            current_dir_template = os.path.join(os.getcwd(), template_filename)
            if os.path.exists(current_dir_template):
                try:
                    # 使用临时文件来预览
                    temp_dir = tempfile.gettempdir()
                    temp_template_name = f"单词表模板_{int(time.time())}.xlsx"
                    temp_template_path = os.path.join(temp_dir, temp_template_name)
                    
                    import shutil
                    shutil.copy2(current_dir_template, temp_template_path)
                    
                    # 使用默认程序打开临时模板文件，并等待程序关闭
                    process = subprocess.Popen(['start', '/wait', temp_template_path], shell=True)
                    # 等待进程结束（即用户关闭表格文件）
                    process.wait()
                    
                    # 关闭后立即删除临时文件
                    if os.path.exists(temp_template_path):
                        try:
                            os.remove(temp_template_path)
                        except:
                            pass
                    
                    self.status_var.set("已打开表格模板")
                    
                except Exception as e:
                    messagebox.showerror("错误", f"打开模板失败: {str(e)}")
                    self.status_var.set("打开模板失败")
            else:
                # 创建一个临时模板文件
                try:
                    from openpyxl import Workbook
                    wb = Workbook()
                    ws = wb.active
                    ws.title = "单词表"
                    headers = ['英文单词', '英文音标', '词根词缀', '例句', '例句释义', '单词释义']
                    for col, header in enumerate(headers, 1):
                        ws.cell(row=1, column=col, value=header)
                    example_data = [
                        ['apple', '/ˈæpl/', 'a-pple', 'I eat an apple every day.', '我每天吃一个苹果。', '苹果'],
                        ['banana', '/bəˈnɑːnə/', 'ban-ana', 'Bananas are yellow.', '香蕉是黄色的。', '香蕉'],
                        ['cat', '/kæt/', 'cat', 'The cat is sleeping.', '猫在睡觉。', '猫'],
                        ['dog', '/dɒɡ/', 'dog', 'Dogs are loyal animals.', '狗是忠诚的动物。', '狗'],
                        ['elephant', '/ˈelɪfənt/', 'ele-ph-ant', 'Elephants have long trunks.', '大象有长长的鼻子。', '大象']
                    ]
                    for row, data in enumerate(example_data, 2):
                        for col, value in enumerate(data, 1):
                            ws.cell(row=row, column=col, value=value)
                    
                    # 创建临时文件而不是在根目录创建
                    temp_dir = tempfile.gettempdir()
                    temp_template_name = f"单词表模板_{int(time.time())}.xlsx"
                    temp_template_path = os.path.join(temp_dir, temp_template_name)
                    
                    # 保存临时模板文件
                    wb.save(temp_template_path)
                    
                    # 打开模板文件，并等待程序关闭
                    process = subprocess.Popen(['start', '/wait', temp_template_path], shell=True)
                    # 等待进程结束（即用户关闭表格文件）
                    process.wait()
                    
                    # 关闭后立即删除临时文件
                    if os.path.exists(temp_template_path):
                        try:
                            os.remove(temp_template_path)
                        except:
                            pass
                    
                    self.status_var.set("已创建并打开表格模板")
                    
                except Exception as e:
                    messagebox.showerror("错误", f"创建模板文件失败: {str(e)}")
                    self.status_var.set("创建模板文件失败")
    
    def load_data_from_sheet(self, input_file: str, sheet_name: str) -> List[Dict[str, str]]:
        """从指定表格加载数据"""
        try:
            wb = load_workbook(input_file)
            ws = wb[sheet_name]
        except Exception as e:
            raise Exception(f"加载表格失败: {str(e)}")
        
        # 获取表头
        headers = [cell.value for cell in ws[1]]
        
        # 验证必要的列是否存在
        required_columns = ['英文单词', '英文音标', '词根词缀', '例句', '例句释义', '单词释义']
        missing_columns = [col for col in required_columns if col not in headers]
        if missing_columns:
            raise Exception(f"表格缺少必要的列: {', '.join(missing_columns)}")
        
        # 加载数据
        data = []
        for row in ws.iter_rows(min_row=2, values_only=True):
            row_dict = {}
            for i, value in enumerate(row):
                if i < len(headers):
                    row_dict[headers[i]] = str(value) if value is not None else ''
            data.append(row_dict)
        
        return data

def load_data(input_file: str) -> List[Dict[str, str]]:
    """从Excel文件加载数据并进行验证"""
    try:
        wb = load_workbook(input_file)
        ws = wb.active
    except Exception as e:
        print(f"读取Excel文件时发生错误：{e}")
        sys.exit(1)
    
    # 获取表头
    headers = [cell.value for cell in ws[1]]
    
    # 验证必要的列是否存在
    missing_columns = [col for col in REQUIRED_COLUMNS if col not in headers]
    if missing_columns:
        print(f"Excel文件缺少必要的列：{', '.join(missing_columns)}")
        sys.exit(1)
    
    # 加载数据
    data = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        row_dict = {}
        for i, value in enumerate(row):
            if i < len(headers):
                row_dict[headers[i]] = str(value) if value is not None else ''
        data.append(row_dict)
    
    return data

def load_data_from_sheet(input_file: str, sheet_name: str) -> List[Dict[str, str]]:
    """从指定表格加载数据"""
    try:
        wb = load_workbook(input_file)
        ws = wb[sheet_name]
    except Exception as e:
        raise Exception(f"加载表格失败: {str(e)}")
    
    # 获取表头
    headers = [cell.value for cell in ws[1]]
    
    # 验证必要的列是否存在
    missing_columns = [col for col in REQUIRED_COLUMNS if col not in headers]
    if missing_columns:
        raise Exception(f"表格缺少必要的列: {', '.join(missing_columns)}")
    
    # 加载数据
    data = []
    for row in ws.iter_rows(min_row=2, values_only=True):
        row_dict = {}
        for i, value in enumerate(row):
            if i < len(headers):
                row_dict[headers[i]] = str(value) if value is not None else ''
        data.append(row_dict)
    
    return data

def generate_ppt_from_sheet(input_file: str, output_file: str, sheet_name: str) -> int:
    """从指定表格生成PPT文件并返回处理的单词数量"""
    # 加载数据
    data = load_data_from_sheet(input_file, sheet_name)
    
    # 创建演示文稿
    ppt = create_presentation()
    
    # 生成幻灯片
    word_count = len(data)
    
    for index, row in enumerate(data):
        generate_slide(ppt, row)
    
    # 保存演示文稿
    try:
        ppt.save(output_file)
        return word_count
    except Exception as e:
        raise Exception(f"保存演示文稿时发生错误: {e}")

def main() -> None:
    """主函数"""
    args = parse_args()
    
    # 如果没有指定命令行参数且支持GUI，则启动图形界面
    if not (args.input or args.output) and HAS_GUI:
        root = tk.Tk()
        app = PPTGeneratorGUI(root)
        root.mainloop()
    else:
        # 使用命令行模式
        generate_ppt(args.input, args.output)

if __name__ == '__main__':
    main()
